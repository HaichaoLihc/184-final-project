#!/usr/bin/env python3
"""Render docs/presentation/index.html to a 1920x1080 PPTX deck.

Approach
--------
The deck is implemented as a `deck-stage` web component that holds N
`<section>` slides and shows one at a time. We:

  1. Boot a local HTTP server rooted at docs/presentation/.
  2. Open the page in headless Chromium at exactly 1920x1080.
  3. Programmatically iterate through every slide via the component's
     `goTo(i)` API, waiting for fonts + lazy images to settle, and PNG-shoot
     each slide.
  4. Assemble the PNGs into a 16:9 PPTX (one slide per PNG, full-bleed).

Run from the repo root:

    python3 scripts/build_presentation_pptx.py
"""
from __future__ import annotations

import asyncio
import http.server
import socketserver
import threading
from contextlib import contextmanager
from pathlib import Path

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
DECK_DIR = ROOT / "docs" / "presentation"
OUT_DIR  = ROOT / "docs" / "presentation" / "exported"
PPTX_OUT = OUT_DIR / "Underwater_Volumetric_Path_Tracing.pptx"

WIDTH  = 1920
HEIGHT = 1080
PORT   = 8765


@contextmanager
def http_server(directory: Path, port: int):
    """Serve `directory` on http://localhost:port for the lifetime of the
    contextmanager. Runs in a daemon thread so the process exits cleanly."""
    handler = lambda *a, **k: http.server.SimpleHTTPRequestHandler(
        *a, directory=str(directory), **k
    )
    httpd = socketserver.ThreadingTCPServer(("127.0.0.1", port), handler)
    httpd.daemon_threads = True
    thread = threading.Thread(target=httpd.serve_forever, daemon=True)
    thread.start()
    try:
        yield f"http://127.0.0.1:{port}"
    finally:
        httpd.shutdown()
        httpd.server_close()


async def render_slides() -> list[Path]:
    """Render every slide section to a PNG. Returns the list of file paths
    in slide order."""
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    out_paths: list[Path] = []
    with http_server(DECK_DIR, PORT) as base:
        url = f"{base}/index.html"
        async with async_playwright() as pw:
            browser = await pw.chromium.launch()
            ctx = await browser.new_context(
                viewport={"width": WIDTH, "height": HEIGHT},
                device_scale_factor=1.0,
            )
            page = await ctx.new_page()
            await page.goto(url, wait_until="networkidle")

            # The deck-stage component lives in a closed-ish shadow root,
            # so a top-level CSS shim can't reach the overlay. Instead, we
            # poke into the shadowRoot once and hide the overlay element by
            # inline style. Also force the page background dark so the area
            # behind the slide stays consistent.
            await page.add_style_tag(content="""
                html, body { background: #07101a !important; }
            """)
            await page.evaluate("""() => {
                const s = document.querySelector('deck-stage');
                if (!s || !s.shadowRoot) return;
                const ov = s.shadowRoot.querySelector('.overlay');
                if (ov) ov.style.display = 'none';
                const tap = s.shadowRoot.querySelector('.tapzones');
                if (tap) tap.style.display = 'none';
            }""")

            # How many slides?
            count = await page.evaluate(
                "document.querySelector('deck-stage').length"
            )
            print(f"[deck] {count} slides")

            for i in range(count):
                # Navigate, then wait for fonts and images on the now-active
                # slide to be fully decoded.
                await page.evaluate(f"document.querySelector('deck-stage').goTo({i})")
                await page.wait_for_timeout(150)
                await page.evaluate("""async () => {
                    if (document.fonts && document.fonts.ready) {
                        await document.fonts.ready;
                    }
                    const imgs = Array.from(document.querySelectorAll(
                        '[data-deck-active] img'
                    ));
                    await Promise.all(imgs.map(img => {
                        if (img.complete && img.naturalWidth > 0) return Promise.resolve();
                        return new Promise((res) => {
                            img.addEventListener('load', res, { once: true });
                            img.addEventListener('error', res, { once: true });
                        });
                    }));
                }""")
                # Brief settle window after the wait (CSS transitions etc.)
                await page.wait_for_timeout(200)

                out_path = OUT_DIR / f"slide_{i+1:02d}.png"
                # Full-page screenshot at the viewport size; deck-stage fills
                # the viewport so this gives us exactly 1920x1080.
                await page.screenshot(path=str(out_path), full_page=False)
                print(f"[deck] slide {i+1:02d} -> {out_path.name}")
                out_paths.append(out_path)

            await browser.close()
    return out_paths


def assemble_pptx(slide_pngs: list[Path]) -> Path:
    """One slide per PNG, full-bleed, 16:9 (1920x1080) deck."""
    prs = Presentation()
    # PowerPoint internal unit is EMU (914400 EMU = 1 inch). Slide layout in
    # inches: 1920px / 96 dpi = 20", 1080px / 96 dpi = 11.25".
    prs.slide_width  = Emu(20 * 914400)
    prs.slide_height = Emu(int(11.25 * 914400))
    blank = prs.slide_layouts[6]  # the "Blank" layout

    for png in slide_pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(png),
            left=0, top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_OUT))
    return PPTX_OUT


def main():
    pngs = asyncio.run(render_slides())
    pptx = assemble_pptx(pngs)
    print(f"[pptx] wrote {pptx} ({pptx.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
